#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
split_anim.py  ―  “アニメーション前後” でスライドを静的に 2 枚へ分割

目的
------
PowerPoint (.pptx / .pptm) を解析し

1. 入口アニメ (appear / fade‑in / click で visible になる など)
2. 出口アニメ (fade‑out / clickEffect hidden / motion path で画面外へ など)

が設定された図形を検出
→ **Before** スライド：入口オブジェクトを非表示
→ **After**  スライド：出口オブジェクトを非表示

LLM 用に「テキストやグラフが隠れない」静的スライドを生成します。

依存
------
    pip install python-pptx lxml
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from zipfile import ZipFile
import posixpath
from typing import Tuple, Set, Dict

from lxml import etree
from pptx import Presentation
from pptx.slide import Slide
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

# ---------- 設定 ---------- #
INPUT = Path(__file__).with_name("input.pptx")
OUTPUT = Path(__file__).with_name("out_split.pptx")

# ---------------------------------------------------------------------
# 1. XML 名前空間辞書  （必要に応じて追加で拡張）
# ---------------------------------------------------------------------
NS: Dict[str, str] = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
}


# ---------------------------------------------------------------------
# 2. アニメーション解析 ― Entrance / Exit 形状 id を検出
# ---------------------------------------------------------------------
def extract_entr_exit_ids(slide_xml: etree._Element) -> Tuple[Set[str], Set[str]]:
    """
    Entrance / Exit の候補となるシェイプ id(spid) を抽出する。

    ■ 判定ロジック
    1. animEffect
       - presetClass="entr"            → Entrance
       - presetClass="exit"            → Exit
       - filter="in:*"                 → Entrance
       - filter="out:*"                → Exit
    2. clickEffect (cTn@nodeType)
       - presetClass="entr"/"exit"     → 同上
       - visibility hidden/visible     → Exit / Entrance
    3. set (opacity / visibility)
       - attrName = style.visibility   + val=hidden/visible → Exit / Entrance
       - attrName = opacity|style.opacity + to='0'/'1'      → Exit / Entrance
    4. animMotion
       - 最終座標がスライド外 (relOff > 100 000 % or < 0 %) → Exit
    5. 3D Arrive/Leave (p14)
       - presetClass="entr"/"exit"     → Entrance / Exit
    6. 予期しない nodeType は無視してログ（ここでは pass）
    """

    def _add(bucket: Set[str], ids: set[str]) -> None:
        bucket.update({i for i in ids if i})

    entrance: Set[str] = set()
    exit_: Set[str] = set()

    # 2‑1) animEffect  (presetClass / filter)
    for eff in slide_xml.xpath(".//p:animEffect", namespaces=NS):
        cls = (eff.get("presetClass") or "").lower()
        filt = (eff.get("filter") or "").lower()
        ids = {t.get("spid") for t in eff.xpath(".//p:spTgt", namespaces=NS)}
        if cls == "entr" or filt.startswith("in:"):
            _add(entrance, ids)
        elif cls == "exit" or filt.startswith("out:"):
            _add(exit_, ids)

    # 2‑2) clickEffect (nodeType) + visibility
    for ctn in slide_xml.xpath(".//p:cTn[@nodeType='clickEffect']", namespaces=NS):
        cls = (ctn.get("presetClass") or "").lower()
        ids = {t.get("spid") for t in ctn.xpath(".//p:spTgt", namespaces=NS)}
        if cls in ("entr", "exit"):
            _add(entrance if cls == "entr" else exit_, ids)
        else:
            vis = ctn.xpath(
                ".//p:set[p:attrName='style.visibility']/p:to/p:strVal/@val",
                namespaces=NS,
            )
            if vis:
                _add(exit_ if vis[0] == "hidden" else entrance, ids)

    # 2‑3) set (visibility / opacity)
    for s in slide_xml.xpath(".//p:set", namespaces=NS):
        ids = {t.get("spid") for t in s.xpath(".//p:spTgt", namespaces=NS)}
        attr = {n.text.lower() for n in s.xpath(".//p:attrName", namespaces=NS)}
        to_val = (
            s.xpath("./p:to/p:strVal/@val", namespaces=NS)
            or s.xpath("./p:to/@val", namespaces=NS)
            or s.xpath("./p:to/@valLst", namespaces=NS)
        )
        to_val = to_val[0] if to_val else ""
        if "style.visibility" in attr:
            _add(exit_ if to_val == "hidden" else entrance, ids)
        elif {"opacity", "style.opacity"} & attr:
            _add(exit_ if to_val in ("0", "0.0") else entrance, ids)

    # 2‑4) motionPath → スライド外
    for mot in slide_xml.xpath(".//p:animMotion", namespaces=NS):
        to_xy = mot.find("./p:to", namespaces=NS)
        if to_xy is not None:
            x = int(to_xy.get("x", "0"))
            y = int(to_xy.get("y", "0"))
            # EMU 100 000 ≒ 100%
            if x < 0 or y < 0 or x > 100000 or y > 100000:
                ids = {t.get("spid") for t in mot.xpath(".//p:spTgt", namespaces=NS)}
                _add(exit_, ids)

    # 2‑5) 3D Arrive / Leave (p14)
    for node in slide_xml.xpath(".//p14:animEffect", namespaces=NS):
        cls = (node.get("presetClass") or "").lower()
        ids = {t.get("spid") for t in node.xpath(".//p:spTgt", namespaces=NS)}
        if cls == "entr":
            _add(entrance, ids)
        elif cls == "exit":
            _add(exit_, ids)

    return entrance, exit_


# ---------------------------------------------------------------------
# 3. スライド複製 ― src の直後にクローンし rels を同期
# ---------------------------------------------------------------------
_RELS_COPY_TYPES = {
    RT.IMAGE,
    RT.CHART,
    RT.MEDIA,
    RT.VIDEO,
    RT.OLE_OBJECT,
}


def clone_slide_after(prs: Presentation, src: Slide) -> Slide:
    """
    - 新しい空白スライドを末尾へ add し
    - sldId を src の直後に移動
    - shapes を deepcopy
    - 画像/チャート等の rels を複製
    Returns 追加した Slide オブジェクト
    """
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)

    # sldId を順序どおり挿入
    sldIdLst = prs.slides._sldIdLst
    src_idx = list(prs.slides).index(src)
    new_id = sldIdLst[-1]
    sldIdLst.remove(new_id)
    sldIdLst.insert(src_idx + 1, new_id)

    # shapes deep copy
    for shp in src.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")

    # rels copy
    for rel in src.part.rels.values():
        if rel.reltype in _RELS_COPY_TYPES and rel.rId not in new_slide.part.rels:
            new_slide.part.relate_to(rel._target, rel.reltype, rel.rId)

    return new_slide


# ---------------------------------------------------------------------
# 4. 図形削除 ― 再帰で親グループごと drop
# ---------------------------------------------------------------------
def drop_shapes(slide: Slide, target_ids: Set[str]) -> None:
    """
    target_ids に含まれる spid を持つ図形
    （およびそれを子に含む grpSp）を削除
    """

    def collect(el: etree._Element, trash: Set[etree._Element]) -> bool:
        """True を返すと親 grpSp も削除対象"""
        hit = el.get("id") in target_ids
        for child in el:
            hit = collect(child, trash) or hit
        if hit:
            trash.add(el)
        return hit

    candidates: Set[etree._Element] = set()
    collect(slide.shapes._spTree, candidates)
    for el in candidates:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


# ---------------------------------------------------------------------
# 5. メイン処理
# ---------------------------------------------------------------------
def split_pptx(in_path: Path, out_path: Path) -> bool:
    prs = Presentation(in_path)
    any_split = False

    # スライド XML パス一覧
    with ZipFile(in_path) as zf:
        slide_parts = sorted(p for p in zf.namelist() if p.startswith("ppt/slides/slide") and p.endswith(".xml"))

    for idx, part in enumerate(slide_parts, start=1):
        with ZipFile(in_path) as zf:
            slide_xml = etree.fromstring(zf.read(part))
        ids_in, ids_out = extract_entr_exit_ids(slide_xml)
        if not ids_in and not ids_out:
            continue

        src_slide = prs.slides[idx - 1]
        # --- Before スライド ---
        before = clone_slide_after(prs, src_slide)
        drop_shapes(before, ids_in)  # 入口だけ消す

        # --- After スライド ---
        after = clone_slide_after(prs, src_slide)
        drop_shapes(after, ids_out)  # 出口を消す

        any_split = True

    prs.save(out_path)
    return any_split


# ---------------------------------------------------------------------
# 6. CLI
# ---------------------------------------------------------------------
if __name__ == "__main__":
    if not INPUT.exists():
        raise SystemExit(f"❌ {INPUT} not found")
    split_pptx(INPUT, OUTPUT)
