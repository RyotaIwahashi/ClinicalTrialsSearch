#!/usr/bin/env python3
"""
pptx_notes_extract.py

A small CLI utility that extracts the speaker-notes text from each slide in a
PowerPoint file (.pptx / .pptm) and prints the result to STDOUT (or saves it
to a JSON file).

Usage:
    python pptx_notes_extract.py /path/to/file.pptx [-o notes.json]

Dependencies:
    pip install python-pptx
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Dict

from pptx import Presentation
from pptx.oxml.ns import qn


def is_slide_hidden(slide) -> bool:
    """
    Slide オブジェクトが「隠しスライド」なら True。
    - <p:sld show="0">  … Hidden
    - show 属性が存在しない or true  … Visible
    """
    # <p:sld> 要素は slide._element
    # qn("p:show") は「名前空間付き show」用のキー
    # "show" は「名前空間なし show」用のキー
    # ファイル生成元による表記ゆれを吸収するため、両方をチェックする
    val = slide._element.get(qn("p:show")) or slide._element.get("show")

    # 非表示スライドのときに出力する
    if val in ("0", "false", "False"):
        print(f"Hidden slide: {slide.slide_id}")
    return val is not None and val in ("0", "false", "False")


def extract_visible_notes(pptx_path: Path) -> dict[int, str]:
    prs = Presentation(pptx_path)
    results = {}
    for idx, s in enumerate(prs.slides, start=1):
        if is_slide_hidden(s):
            continue  # 非表示スライドをスキップ
        note = s.notes_slide.notes_text_frame.text.strip() if s.has_notes_slide else ""
        results[idx] = note
    return results


def main() -> None:
    # parser = argparse.ArgumentParser(description="Extract speaker notes from a PowerPoint file.")
    # parser.add_argument("pptx", type=Path, help="Path to .pptx or .pptm file")
    # parser.add_argument(
    #     "-o",
    #     "--output",
    #     type=Path,
    #     help="Write results to JSON instead of printing to STDOUT",
    # )
    # args = parser.parse_args()

    # if not args.pptx.exists():
    #     parser.error(f"File not found: {args.pptx}")

    # notes_map = extract_visible_notes(args.pptx)

    # if args.output:
    #     args.output.write_text(json.dumps(notes_map, ensure_ascii=False, indent=2))
    #     print(f"✓ Notes written to {args.output.resolve()}")
    # else:
    #     # Pretty-print to console (one slide per block)
    #     for idx in sorted(notes_map):
    #         print(f"\n=== Slide {idx} ===")
    #         print(notes_map[idx] or "[ no notes ]")

    notes_map = extract_visible_notes("hide_input.pptx")
    # Pretty-print to console (one slide per block)
    for idx in sorted(notes_map):
        print(f"\n=== Slide {idx} ===")
        print(notes_map[idx] or "[ no notes ]")


if __name__ == "__main__":
    main()
