#!/usr/bin/env python3
"""
pptx_split_anim_inout.py
-------------------
* input.pptx を読み込み
* entrance / exit アニメ付きシェイプを検出
* スライドを「Before（入口前）」＋「After（出口後）」の 2 枚に静的分割
* 出力ファイルを out_split.pptx として保存

依存:
    pip install python-pptx lxml
"""

from copy import deepcopy
from pathlib import Path, PurePosixPath
from zipfile import ZipFile
import posixpath
from typing import Set, Tuple

from lxml import etree
from pptx import Presentation
from pptx.slide import Slide  # ← これを既存の import 群に追加

# ---------- 設定 ---------- #
INPUT = Path(__file__).with_name("input.pptx")
OUTPUT = Path(__file__).with_name("out_split.pptx")

# PML = {
#     "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
#     "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
#     "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
# }
PML = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "a14": "http://schemas.microsoft.com/office/drawing/2010/main",
}


# ---------- ユーティリティ ---------- #
def _xml(zf: ZipFile, part: str) -> etree._Element:
    return etree.fromstring(zf.read(part))


def _slide_xml_parts(zf: ZipFile) -> list[str]:
    """ppt/slides/slideX.xml の ZIP 内パスをソートして返す"""
    return sorted(p for p in zf.namelist() if p.startswith("ppt/slides/slide") and p.endswith(".xml"))


# def _entrance_exit_ids(slide_xml: etree._Element) -> Tuple[Set[str], Set[str]]:
#     """
#     <p:animEffect presetClass="entr|exit"> を走査し，
#     spTgt/@spid (shape id) をセットで返す
#     """
#     entr, exit_ = set(), set()
#     for anim in slide_xml.xpath(".//p:animEffect", namespaces=PML):
#         cls = anim.get("presetClass")  # 'entr', 'exit', 'emph' など
#         spids = {tgt.get("spid") for tgt in anim.xpath(".//p:spTgt", namespaces=PML)}
#         if cls == "entr":
#             entr.update(spids)
#         elif cls == "exit":
#             exit_.update(spids)
#     return entr, exit_


def _extract_entr_exit_ids(slide_xml: etree._Element) -> Tuple[Set[str], Set[str]]:
    """
    Entrance / Exit の候補となるシェイプ id(spid) を抽出する。

    ■ 判定ロジック
    1. animEffect
       - presetClass="entr"            → Entrance
       - presetClass="exit"            → Exit
       - filter="in:*"                 → Entrance
       - filter="out:*"                → Exit
    2. clickEffect (cTn@nodeType)
       - presetClass="entr"/"exit"     → 同上
       - visibility hidden/visible     → Exit / Entrance
    3. set (opacity / visibility)
       - attrName = style.visibility   + val=hidden/visible → Exit / Entrance
       - attrName = opacity|style.opacity + to='0'/'1'      → Exit / Entrance
    4. animMotion
       - 最終座標がスライド外 (relOff > 100 000 % or < 0 %) → Exit
    5. 3D Arrive/Leave (p14)
       - presetClass="entr"/"exit"     → Entrance / Exit
    6. 予期しない nodeType は無視してログ（ここでは pass）
    """
    entr, exit_ = set(), set()

    def _add(ids: Set[str], target: str):
        if target == "entr":
            entr.update(ids)
        elif target == "exit":
            exit_.update(ids)

    # ￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣ #
    # 1. <p:animEffect>                                           #
    # ＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿ #
    for node in slide_xml.xpath(".//p:animEffect", namespaces=PML):
        cls = (node.get("presetClass") or "").lower()
        filt = (node.get("filter") or "").lower()
        ids = {tgt.get("spid") for tgt in node.xpath(".//p:spTgt", namespaces=PML)}

        if cls == "entr" or filt.startswith("in:"):
            _add(ids, "entr")
        elif cls == "exit" or filt.startswith("out:"):
            _add(ids, "exit")

    # ￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣ #
    # 2. clickEffect                                              #
    # ＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿ #
    for ctn in slide_xml.xpath(".//p:cTn[@nodeType='clickEffect']", namespaces=PML):
        cls = (ctn.get("presetClass") or "").lower()
        ids = {t.get("spid") for t in ctn.xpath(".//p:spTgt", namespaces=PML)}

        if cls in ("entr", "exit"):
            _add(ids, cls)
            continue

        # visibility が hidden → Exit / visible → Entrance
        vis = ctn.xpath(
            ".//p:set[p:attrName='style.visibility']/p:to/p:strVal/@val",
            namespaces=PML,
        )
        if vis:
            _add(ids, "exit" if vis[0] == "hidden" else "entr")

    # ￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣ #
    # 3. <p:set> で visibility or opacity を変更                #
    # ＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿ #
    for s in slide_xml.xpath(".//p:set", namespaces=PML):
        attr_names = {n.text.lower() for n in s.xpath(".//p:attrName", namespaces=PML)}
        ids = {t.get("spid") for t in s.xpath(".//p:spTgt", namespaces=PML)}
        to = (
            s.xpath("./p:to/p:strVal/@val", namespaces=PML)
            or s.xpath("./p:to/@val", namespaces=PML)
            or s.xpath("./p:to/@valLst", namespaces=PML)
        )
        to_val = to[0] if to else ""

        # visibility
        if "style.visibility" in attr_names:
            _add(ids, "exit" if to_val == "hidden" else "entr")
            continue

        # opacity
        if "opacity" in attr_names or "style.opacity" in attr_names:
            _add(ids, "exit" if to_val in ("0", "0.0") else "entr")

    # ￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣ #
    # 4. Motion Path がキャンバス外へ飛び出す                  #
    # ＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿ #
    for mot in slide_xml.xpath(".//p:animMotion", namespaces=PML):
        ids = {t.get("spid") for t in mot.xpath(".//p:spTgt", namespaces=PML)}
        # <p:to x="200000" y="0" /> などの場合はスライド外（EMU で 100 000 = 幅 100%）
        to_xy = mot.xpath("./p:to", namespaces=PML)
        if to_xy:
            x = int(to_xy[0].get("x", "0"))
            y = int(to_xy[0].get("y", "0"))
            if x < 0 or y < 0 or x > 100000 or y > 100000:
                _add(ids, "exit")

    # ￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣ #
    # 5. 3D モデル Arrive / Leave (p14 名前空間)                #
    # ＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿ #
    for node in slide_xml.xpath(".//p14:animEffect", namespaces=PML):
        cls = (node.get("presetClass") or "").lower()
        ids = {t.get("spid") for t in node.xpath(".//p:spTgt", namespaces=PML)}
        if cls == "entr":
            _add(ids, "entr")
        elif cls == "exit":
            _add(ids, "exit")

    # ￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣￣ #
    # 重複は set で自動解決。戻り値は (entrance_ids, exit_ids) #
    # ＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿＿ #
    return entr, exit_


# def _clone_slide(prs: Presentation, src_slide) -> Slide:
#     """python-pptx に clone API が無いので XML 丸ごと deep copy"""
#     blank = prs.slide_layouts[6]  # 空白レイアウト
#     new_slide = prs.slides.add_slide(blank)
#     # deep copy 元 XML を置換
#     new_slide._element.getparent().replace(new_slide._element, deepcopy(src_slide._element))
#     return new_slide


# def _clone_slide(prs: Presentation, src: Slide) -> Slide:
#     """
#     - 空白レイアウトで新スライドを作成
#     - src.shapes の XML ノードを deep copy して挿入
#     - ノートやタイムラインはコピーしない（必要なら追加実装）
#     """
#     blank = prs.slide_layouts[6]  # 6 = Title Only / Blank など空白
#     dst = prs.slides.add_slide(blank)

#     # 背景やマスター依存スタイルが要る場合はここでコピーする
#     #   dst.background = deepcopy(src.background)  # ←背景色だけならこの行でOK
#     #
#     # Shapes を丸ごと複製
#     for shape in src.shapes:
#         new_el = deepcopy(shape.element)
#         # extLst の直前 (= spTree の末尾) に挿入
#         dst.shapes._spTree.insert_element_before(new_el, "p:extLst")

#     return dst


def _clone_slide(prs, src):
    """src の直後にクローンスライドを挿入して返す"""
    # ① 空白スライドを末尾へ追加
    blank_layout = prs.slide_layouts[6]  # 空白レイアウト
    new_slide = prs.slides.add_slide(blank_layout)

    # ② sldId を src の直後へ移動
    sldIdLst = prs.slides._sldIdLst  # <p:sldIdLst>
    src_idx = list(prs.slides).index(src)  # 0-origin
    new_id = sldIdLst[-1]  # 追加した sldId
    sldIdLst.remove(new_id)
    sldIdLst.insert(src_idx + 1, new_id)

    # ③ src.shapes を deep‑copy
    for shp in src.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shp.element), "p:extLst")

    # ④ 画像・グラフなどの rel をコピー
    for rel in src.part.rels:
        if rel.reltype in (
            # 画像／チャート／メディアなど埋め込み系だけ
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
        ):
            if rel.rId in new_slide.part.rels:
                continue
            new_slide.part.relate_to(rel._target, rel.reltype, rel.rId)

    return new_slide


def _drop_shapes(slide, ids: Set[str]) -> None:
    """shape.id が ids に含まれる要素を削除"""
    for shp in list(slide.shapes):  # list() でコピーしてからでないと iterator 崩れる
        if str(shp.element.get("id")) in ids:
            shp.element.getparent().remove(shp.element)


# ---------- メイン ---------- #
def split_animation_slides(in_path: Path, out_path: Path) -> None:
    prs = Presentation(in_path)  # python-pptx オブジェクト
    with ZipFile(in_path) as zf:  # XML 低レイヤ解析
        slide_parts = _slide_xml_parts(zf)

        # enumerate は 0-origin なので +1 が pptx の SlideIndex
        for idx, slide_part in enumerate(slide_parts, start=1):
            xml_root = _xml(zf, slide_part)
            entr_ids, exit_ids = _extract_entr_exit_ids(xml_root)

            # アニメなしならスキップ
            if not (entr_ids or exit_ids):
                continue

            src_slide = prs.slides[idx - 1]

            before_slide = _clone_slide(prs, src_slide)
            after_slide = _clone_slide(prs, before_slide)

            # --- Before（入口前）: entrance シェイプを除去 --- #
            # --- After（出口後）: exit シェイプを除去 --- #
            _drop_shapes(before_slide, entr_ids)
            _drop_shapes(after_slide, exit_ids)

            # (Optional) どちらにも ' (before)' / ' (after)' タイトルを付けたい場合
            # if before_slide.shapes.title:
            #     before_slide.shapes.title.text += " (before)"
            # if after_slide.shapes.title:
            #     after_slide.shapes.title.text += " (after)"

    prs.save(out_path)
    print(f"✅ 分割完了: {out_path.absolute()}")


# ---------- 実行部 ---------- #
if __name__ == "__main__":
    if not INPUT.exists():
        raise SystemExit(f"❌ {INPUT} not found")
    split_animation_slides(INPUT, OUTPUT)
